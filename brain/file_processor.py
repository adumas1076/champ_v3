# ============================================
# CHAMP V3 — File Processor
# Brick 10: Universal Format Processing
# Takes any file, detects type, extracts content.
# ============================================
# "Built to build. Born to create."

import base64
import csv
import io
import json
import logging
import mimetypes
import os
import zipfile
import tarfile
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Supported MIME type -> handler mapping
SUPPORTED_TYPES = {
    # Documents
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/rtf": "rtf",
    "text/rtf": "rtf",
    "application/vnd.oasis.opendocument.text": "odt",
    "application/vnd.oasis.opendocument.spreadsheet": "ods",
    "application/vnd.oasis.opendocument.presentation": "odp",
    "text/markdown": "text",
    "text/plain": "text",
    # Spreadsheets
    "text/csv": "csv",
    "text/tab-separated-values": "csv",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    # Data
    "application/json": "json_file",
    "application/xml": "xml",
    "text/xml": "xml",
    "application/x-yaml": "yaml",
    "text/yaml": "yaml",
    "application/toml": "toml_file",
    "application/vnd.apache.parquet": "parquet",
    "application/x-sqlite3": "sqlite",
    # Images
    "image/png": "image",
    "image/jpeg": "image",
    "image/gif": "image",
    "image/webp": "image",
    "image/bmp": "image",
    "image/tiff": "image",
    "image/svg+xml": "svg",
    # Audio
    "audio/mpeg": "audio",
    "audio/wav": "audio",
    "audio/x-wav": "audio",
    "audio/mp4": "audio",
    "audio/m4a": "audio",
    "audio/ogg": "audio",
    "audio/flac": "audio",
    "audio/aac": "audio",
    # Video
    "video/mp4": "video",
    "video/quicktime": "video",
    "video/webm": "video",
    "video/x-msvideo": "video",
    "video/x-matroska": "video",
    # Archives
    "application/zip": "zip",
    "application/x-tar": "tar",
    "application/gzip": "tar",
    "application/x-rar-compressed": "rar",
    # Email
    "message/rfc822": "eml",
    "application/vnd.ms-outlook": "msg",
    # Web
    "text/html": "html",
    # Code (all treated as text)
    "text/x-python": "text",
    "application/javascript": "text",
    "text/javascript": "text",
    "text/css": "text",
    "application/sql": "text",
    "application/typescript": "text",
}

# Extension fallbacks for when MIME detection fails
EXT_FALLBACKS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".rtf": "rtf",
    ".odt": "odt",
    ".ods": "ods",
    ".odp": "odp",
    ".csv": "csv",
    ".tsv": "csv",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".json": "json_file",
    ".xml": "xml",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".toml": "toml_file",
    ".parquet": "parquet",
    ".db": "sqlite",
    ".sqlite": "sqlite",
    ".sqlite3": "sqlite",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".webp": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
    ".svg": "svg",
    ".mp3": "audio",
    ".wav": "audio",
    ".m4a": "audio",
    ".ogg": "audio",
    ".flac": "audio",
    ".aac": "audio",
    ".mp4": "video",
    ".mov": "video",
    ".webm": "video",
    ".avi": "video",
    ".mkv": "video",
    ".zip": "zip",
    ".tar": "tar",
    ".gz": "tar",
    ".tar.gz": "tar",
    ".bz2": "tar",
    ".rar": "rar",
    ".eml": "eml",
    ".msg": "msg",
    ".html": "html",
    ".htm": "html",
    ".py": "text",
    ".js": "text",
    ".ts": "text",
    ".tsx": "text",
    ".jsx": "text",
    ".css": "text",
    ".sql": "text",
    ".md": "text",
    ".txt": "text",
    ".log": "text",
    ".ini": "text",
    ".cfg": "text",
    ".env": "text",
    ".sh": "text",
    ".bash": "text",
    ".ps1": "text",
    ".bat": "text",
    ".rs": "text",
    ".go": "text",
    ".java": "text",
    ".c": "text",
    ".cpp": "text",
    ".h": "text",
    ".rb": "text",
    ".php": "text",
    ".swift": "text",
    ".kt": "text",
    ".r": "text",
    ".lua": "text",
}


class FileProcessResult:
    """Result of processing a file."""

    def __init__(
        self,
        text: str,
        file_type: str,
        filename: str,
        metadata: Optional[dict] = None,
        image_b64: Optional[str] = None,
        mime_type: Optional[str] = None,
    ):
        self.text = text
        self.file_type = file_type
        self.filename = filename
        self.metadata = metadata or {}
        self.image_b64 = image_b64  # For images — base64 encoded for vision models
        self.mime_type = mime_type

    def to_dict(self) -> dict:
        d = {
            "text": self.text,
            "file_type": self.file_type,
            "filename": self.filename,
            "metadata": self.metadata,
        }
        if self.image_b64:
            d["image_b64"] = self.image_b64
        if self.mime_type:
            d["mime_type"] = self.mime_type
        return d


def detect_handler(filename: str, content_type: Optional[str] = None) -> str:
    """Detect which handler to use for a file."""
    # Try MIME type first
    if content_type and content_type in SUPPORTED_TYPES:
        return SUPPORTED_TYPES[content_type]

    # Fall back to extension
    ext = Path(filename).suffix.lower()
    if ext in EXT_FALLBACKS:
        return EXT_FALLBACKS[ext]

    # Try mimetypes module
    guessed, _ = mimetypes.guess_type(filename)
    if guessed and guessed in SUPPORTED_TYPES:
        return SUPPORTED_TYPES[guessed]

    return "unknown"


async def process_file(
    file_bytes: bytes,
    filename: str,
    content_type: Optional[str] = None,
) -> FileProcessResult:
    """
    Process any file and extract its content.
    Returns a FileProcessResult with extracted text and metadata.
    """
    handler = detect_handler(filename, content_type)
    logger.info(f"[FILE] Processing {filename} ({len(file_bytes)} bytes) -> handler: {handler}")

    try:
        if handler == "pdf":
            return _process_pdf(file_bytes, filename)
        elif handler == "docx":
            return _process_docx(file_bytes, filename)
        elif handler == "pptx":
            return _process_pptx(file_bytes, filename)
        elif handler == "rtf":
            return _process_rtf(file_bytes, filename)
        elif handler == "odt":
            return _process_odt(file_bytes, filename)
        elif handler == "csv":
            return _process_csv(file_bytes, filename)
        elif handler == "xlsx":
            return _process_xlsx(file_bytes, filename)
        elif handler == "json_file":
            return _process_json(file_bytes, filename)
        elif handler == "xml":
            return _process_xml(file_bytes, filename)
        elif handler == "yaml":
            return _process_yaml(file_bytes, filename)
        elif handler == "toml_file":
            return _process_toml(file_bytes, filename)
        elif handler == "parquet":
            return _process_parquet(file_bytes, filename)
        elif handler == "sqlite":
            return _process_sqlite(file_bytes, filename)
        elif handler == "image":
            return _process_image(file_bytes, filename, content_type)
        elif handler == "svg":
            return _process_svg(file_bytes, filename)
        elif handler == "audio":
            return _process_audio(file_bytes, filename, content_type)
        elif handler == "video":
            return _process_video(file_bytes, filename)
        elif handler == "zip":
            return _process_zip(file_bytes, filename)
        elif handler == "tar":
            return _process_tar(file_bytes, filename)
        elif handler == "rar":
            return _process_rar(file_bytes, filename)
        elif handler == "eml":
            return _process_eml(file_bytes, filename)
        elif handler == "msg":
            return _process_msg(file_bytes, filename)
        elif handler == "html":
            return _process_html(file_bytes, filename)
        elif handler == "text":
            return _process_text(file_bytes, filename)
        else:
            # Try as text, fall back to binary summary
            try:
                text = file_bytes.decode("utf-8")
                return FileProcessResult(
                    text=text[:50000],
                    file_type="text",
                    filename=filename,
                    metadata={"note": "Unknown type, treated as text"},
                )
            except UnicodeDecodeError:
                return FileProcessResult(
                    text=f"[Binary file: {filename}, {len(file_bytes)} bytes. Cannot extract text.]",
                    file_type="binary",
                    filename=filename,
                )
    except Exception as e:
        logger.error(f"[FILE] Failed to process {filename}: {e}")
        return FileProcessResult(
            text=f"[Error processing {filename}: {e}]",
            file_type="error",
            filename=filename,
            metadata={"error": str(e)},
        )


# ---- Document Handlers ----

def _process_pdf(data: bytes, filename: str) -> FileProcessResult:
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text}")
    doc.close()

    return FileProcessResult(
        text="\n\n".join(pages) if pages else "[PDF has no extractable text]",
        file_type="pdf",
        filename=filename,
        metadata={"pages": len(pages)},
    )


def _process_docx(data: bytes, filename: str) -> FileProcessResult:
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract tables
    tables_text = []
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        tables_text.append(f"--- Table {i + 1} ---\n" + "\n".join(rows))

    full_text = "\n\n".join(paragraphs)
    if tables_text:
        full_text += "\n\n" + "\n\n".join(tables_text)

    return FileProcessResult(
        text=full_text if full_text.strip() else "[DOCX has no extractable text]",
        file_type="docx",
        filename=filename,
        metadata={"paragraphs": len(paragraphs), "tables": len(doc.tables)},
    )


def _process_pptx(data: bytes, filename: str) -> FileProcessResult:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))

    return FileProcessResult(
        text="\n\n".join(slides) if slides else "[PPTX has no extractable text]",
        file_type="pptx",
        filename=filename,
        metadata={"slides": len(prs.slides)},
    )


def _process_rtf(data: bytes, filename: str) -> FileProcessResult:
    from striprtf.striprtf import rtf_to_text

    text = rtf_to_text(data.decode("utf-8", errors="replace"))
    return FileProcessResult(text=text, file_type="rtf", filename=filename)


def _process_odt(data: bytes, filename: str) -> FileProcessResult:
    from odf.opendocument import load as odf_load
    from odf.text import P

    doc = odf_load(io.BytesIO(data))
    paragraphs = []
    for p in doc.getElementsByType(P):
        text = ""
        for node in p.childNodes:
            if hasattr(node, "data"):
                text += node.data
            elif hasattr(node, "__str__"):
                text += str(node)
        if text.strip():
            paragraphs.append(text)

    return FileProcessResult(
        text="\n\n".join(paragraphs) if paragraphs else "[ODT has no extractable text]",
        file_type="odt",
        filename=filename,
        metadata={"paragraphs": len(paragraphs)},
    )


# ---- Spreadsheet & Data Handlers ----

def _process_csv(data: bytes, filename: str) -> FileProcessResult:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)

    # Format as readable table
    formatted = []
    for i, row in enumerate(rows[:200]):  # Cap at 200 rows
        formatted.append(" | ".join(row))

    return FileProcessResult(
        text="\n".join(formatted),
        file_type="csv",
        filename=filename,
        metadata={"rows": len(rows), "displayed": min(len(rows), 200)},
    )


def _process_xlsx(data: bytes, filename: str) -> FileProcessResult:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
    wb.close()

    return FileProcessResult(
        text="\n\n".join(sheets) if sheets else "[XLSX has no data]",
        file_type="xlsx",
        filename=filename,
        metadata={"sheets": len(wb.sheetnames)},
    )


def _process_json(data: bytes, filename: str) -> FileProcessResult:
    parsed = json.loads(data.decode("utf-8"))
    text = json.dumps(parsed, indent=2)
    if len(text) > 50000:
        text = text[:50000] + "\n... [truncated]"
    return FileProcessResult(text=text, file_type="json", filename=filename)


def _process_xml(data: bytes, filename: str) -> FileProcessResult:
    from lxml import etree

    tree = etree.fromstring(data)
    text = etree.tostring(tree, pretty_print=True, encoding="unicode")
    if len(text) > 50000:
        text = text[:50000] + "\n... [truncated]"
    return FileProcessResult(text=text, file_type="xml", filename=filename)


def _process_yaml(data: bytes, filename: str) -> FileProcessResult:
    import yaml

    text = data.decode("utf-8")
    # Validate it's valid YAML
    yaml.safe_load(text)
    return FileProcessResult(text=text, file_type="yaml", filename=filename)


def _process_toml(data: bytes, filename: str) -> FileProcessResult:
    import toml as toml_lib

    parsed = toml_lib.loads(data.decode("utf-8"))
    text = json.dumps(parsed, indent=2)
    return FileProcessResult(text=text, file_type="toml", filename=filename)


def _process_parquet(data: bytes, filename: str) -> FileProcessResult:
    import pyarrow.parquet as pq

    table = pq.read_table(io.BytesIO(data))
    # Convert first 200 rows to readable format
    df_dict = table.slice(0, 200).to_pydict()
    headers = list(df_dict.keys())
    rows_count = len(next(iter(df_dict.values()))) if df_dict else 0

    lines = [" | ".join(headers)]
    for i in range(rows_count):
        row = [str(df_dict[h][i]) for h in headers]
        lines.append(" | ".join(row))

    return FileProcessResult(
        text="\n".join(lines),
        file_type="parquet",
        filename=filename,
        metadata={"columns": len(headers), "total_rows": table.num_rows},
    )


def _process_sqlite(data: bytes, filename: str) -> FileProcessResult:
    import sqlite3
    import tempfile

    # SQLite needs a file, write temp
    with tempfile.NamedTemporaryFile(suffix=".db", delete=False) as f:
        f.write(data)
        tmp_path = f.name

    try:
        conn = sqlite3.connect(tmp_path)
        cursor = conn.cursor()

        # Get tables
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
        tables = [r[0] for r in cursor.fetchall()]

        parts = []
        for table in tables:
            cursor.execute(f"SELECT * FROM [{table}] LIMIT 50")
            rows = cursor.fetchall()
            cols = [d[0] for d in cursor.description]
            header = " | ".join(cols)
            data_rows = [" | ".join(str(c) for c in row) for row in rows]
            parts.append(f"--- Table: {table} ---\n{header}\n" + "\n".join(data_rows))

        conn.close()
        return FileProcessResult(
            text="\n\n".join(parts) if parts else "[Empty database]",
            file_type="sqlite",
            filename=filename,
            metadata={"tables": tables},
        )
    finally:
        os.unlink(tmp_path)


# ---- Image Handlers ----

def _process_image(data: bytes, filename: str, content_type: Optional[str] = None) -> FileProcessResult:
    from PIL import Image

    img = Image.open(io.BytesIO(data))
    width, height = img.size
    mode = img.mode
    fmt = img.format

    # Encode as base64 for vision model
    b64 = base64.b64encode(data).decode("utf-8")

    # Detect MIME for the image_url format
    mime = content_type or f"image/{(fmt or 'png').lower()}"

    return FileProcessResult(
        text=f"[Image: {filename} | {width}x{height} | {mode} | {fmt}]",
        file_type="image",
        filename=filename,
        metadata={"width": width, "height": height, "mode": mode, "format": fmt},
        image_b64=b64,
        mime_type=mime,
    )


def _process_svg(data: bytes, filename: str) -> FileProcessResult:
    # SVG is XML — extract as text
    text = data.decode("utf-8", errors="replace")
    return FileProcessResult(
        text=text if len(text) < 50000 else text[:50000] + "\n... [truncated]",
        file_type="svg",
        filename=filename,
        metadata={"note": "SVG read as XML source. Visual rendering requires libcairo."},
    )


# ---- Audio Handlers ----

def _process_audio(data: bytes, filename: str, content_type: Optional[str] = None) -> FileProcessResult:
    from pydub import AudioSegment

    ext = Path(filename).suffix.lower().lstrip(".")
    if ext == "m4a":
        ext = "mp4"

    try:
        audio = AudioSegment.from_file(io.BytesIO(data), format=ext)
        duration_sec = len(audio) / 1000.0
        channels = audio.channels
        sample_rate = audio.frame_rate

        return FileProcessResult(
            text=(
                f"[Audio file: {filename} | Duration: {duration_sec:.1f}s | "
                f"Channels: {channels} | Sample rate: {sample_rate}Hz]\n\n"
                "To transcribe this audio, use the Whisper API or send to the voice pipeline."
            ),
            file_type="audio",
            filename=filename,
            metadata={
                "duration_seconds": duration_sec,
                "channels": channels,
                "sample_rate": sample_rate,
            },
        )
    except Exception as e:
        return FileProcessResult(
            text=f"[Audio file: {filename} | Could not read metadata: {e}]",
            file_type="audio",
            filename=filename,
            metadata={"error": str(e)},
        )


# ---- Video Handlers ----

def _process_video(data: bytes, filename: str) -> FileProcessResult:
    import tempfile
    try:
        from moviepy import VideoFileClip
    except ImportError:
        from moviepy.editor import VideoFileClip

    with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as f:
        f.write(data)
        tmp_path = f.name

    try:
        clip = VideoFileClip(tmp_path)
        duration = clip.duration
        size = clip.size  # (width, height)
        fps = clip.fps
        clip.close()

        return FileProcessResult(
            text=(
                f"[Video file: {filename} | Duration: {duration:.1f}s | "
                f"Resolution: {size[0]}x{size[1]} | FPS: {fps}]\n\n"
                "To analyze video content, extract frames and send to Gemini Flash (vision model)."
            ),
            file_type="video",
            filename=filename,
            metadata={
                "duration_seconds": duration,
                "width": size[0],
                "height": size[1],
                "fps": fps,
            },
        )
    except Exception as e:
        return FileProcessResult(
            text=f"[Video file: {filename} | Could not read metadata: {e}]",
            file_type="video",
            filename=filename,
            metadata={"error": str(e)},
        )
    finally:
        os.unlink(tmp_path)


# ---- Archive Handlers ----

def _process_zip(data: bytes, filename: str) -> FileProcessResult:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        file_list = zf.namelist()
        sizes = {name: zf.getinfo(name).file_size for name in file_list}

    listing = "\n".join(
        f"  {name} ({sizes[name]:,} bytes)" for name in file_list[:100]
    )
    return FileProcessResult(
        text=f"ZIP archive: {filename}\nFiles ({len(file_list)}):\n{listing}",
        file_type="zip",
        filename=filename,
        metadata={"file_count": len(file_list), "files": file_list[:100]},
    )


def _process_tar(data: bytes, filename: str) -> FileProcessResult:
    with tarfile.open(fileobj=io.BytesIO(data)) as tf:
        members = tf.getmembers()
        listing = "\n".join(
            f"  {m.name} ({m.size:,} bytes)" for m in members[:100]
        )

    return FileProcessResult(
        text=f"TAR archive: {filename}\nFiles ({len(members)}):\n{listing}",
        file_type="tar",
        filename=filename,
        metadata={"file_count": len(members)},
    )


def _process_rar(data: bytes, filename: str) -> FileProcessResult:
    import tempfile
    import rarfile

    with tempfile.NamedTemporaryFile(suffix=".rar", delete=False) as f:
        f.write(data)
        tmp_path = f.name

    try:
        with rarfile.RarFile(tmp_path) as rf:
            file_list = rf.namelist()
            listing = "\n".join(f"  {name}" for name in file_list[:100])

        return FileProcessResult(
            text=f"RAR archive: {filename}\nFiles ({len(file_list)}):\n{listing}",
            file_type="rar",
            filename=filename,
            metadata={"file_count": len(file_list)},
        )
    finally:
        os.unlink(tmp_path)


# ---- Email Handlers ----

def _process_eml(data: bytes, filename: str) -> FileProcessResult:
    import email
    from email import policy

    msg = email.message_from_bytes(data, policy=policy.default)
    subject = msg.get("subject", "(no subject)")
    sender = msg.get("from", "(unknown)")
    to = msg.get("to", "(unknown)")
    date = msg.get("date", "(unknown)")

    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_content()
                break
    else:
        body = msg.get_content()

    text = (
        f"From: {sender}\n"
        f"To: {to}\n"
        f"Date: {date}\n"
        f"Subject: {subject}\n"
        f"---\n{body}"
    )

    return FileProcessResult(
        text=text,
        file_type="eml",
        filename=filename,
        metadata={"subject": subject, "from": sender, "to": to},
    )


def _process_msg(data: bytes, filename: str) -> FileProcessResult:
    import tempfile
    import extract_msg

    with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as f:
        f.write(data)
        tmp_path = f.name

    try:
        msg = extract_msg.Message(tmp_path)
        text = (
            f"From: {msg.sender}\n"
            f"To: {msg.to}\n"
            f"Date: {msg.date}\n"
            f"Subject: {msg.subject}\n"
            f"---\n{msg.body}"
        )
        msg.close()

        return FileProcessResult(
            text=text,
            file_type="msg",
            filename=filename,
            metadata={"subject": msg.subject},
        )
    finally:
        os.unlink(tmp_path)


# ---- Web Handler ----

def _process_html(data: bytes, filename: str) -> FileProcessResult:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "html.parser")

    # Remove script and style elements
    for element in soup(["script", "style"]):
        element.decompose()

    text = soup.get_text(separator="\n", strip=True)
    title = soup.title.string if soup.title else None

    return FileProcessResult(
        text=text[:50000] if text else "[HTML has no extractable text]",
        file_type="html",
        filename=filename,
        metadata={"title": title},
    )


# ---- Text Handler ----

def _process_text(data: bytes, filename: str) -> FileProcessResult:
    text = data.decode("utf-8", errors="replace")
    ext = Path(filename).suffix.lower()
    lang = {
        ".py": "python", ".js": "javascript", ".ts": "typescript",
        ".tsx": "typescript", ".jsx": "javascript", ".css": "css",
        ".sql": "sql", ".sh": "bash", ".rs": "rust", ".go": "go",
        ".java": "java", ".c": "c", ".cpp": "cpp", ".rb": "ruby",
        ".php": "php", ".swift": "swift", ".kt": "kotlin",
    }.get(ext, "")

    if lang:
        text = f"```{lang}\n{text}\n```"

    return FileProcessResult(
        text=text[:50000] if len(text) > 50000 else text,
        file_type="text",
        filename=filename,
        metadata={"language": lang} if lang else {},
    )
