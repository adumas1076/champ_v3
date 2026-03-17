# ============================================
# CHAMP V3 — Gate Test: File Processor
# Tests all 30+ format handlers end-to-end
# ============================================

import asyncio
import io
import json
import os
import sys
import struct
import tempfile
import time
import zipfile

# Fix Windows console encoding
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
sys.stderr.reconfigure(encoding='utf-8', errors='replace')

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from brain.file_processor import detect_handler, process_file, FileProcessResult

# ---- Test Helpers ----

PASS = 0
FAIL = 0
SKIP = 0

def result(name, passed, detail=""):
    global PASS, FAIL
    if passed:
        PASS += 1
        print(f"  ✅ {name}{f' — {detail}' if detail else ''}")
    else:
        FAIL += 1
        print(f"  ❌ {name}{f' — {detail}' if detail else ''}")

def skip(name, reason):
    global SKIP
    SKIP += 1
    print(f"  ⏭️  {name} — SKIPPED: {reason}")


# ============================================
# PART 1: Handler Detection (all extensions)
# ============================================

def test_handler_detection():
    print("\n=== PART 1: Handler Detection ===")

    cases = {
        # Documents
        "report.pdf": "pdf",
        "doc.docx": "docx",
        "old.doc": "docx",
        "slides.pptx": "pptx",
        "old_slides.ppt": "pptx",
        "letter.rtf": "rtf",
        "libre.odt": "odt",
        "libre_sheet.ods": "ods",
        "libre_pres.odp": "odp",
        "readme.md": "text",
        "notes.txt": "text",
        # Spreadsheets
        "data.csv": "csv",
        "data.tsv": "csv",
        "workbook.xlsx": "xlsx",
        "old_workbook.xls": "xlsx",
        # Data/Config
        "config.json": "json_file",
        "feed.xml": "xml",
        "config.yaml": "yaml",
        "config.yml": "yaml",
        "settings.toml": "toml_file",
        "data.parquet": "parquet",
        "app.db": "sqlite",
        "app.sqlite": "sqlite",
        "app.sqlite3": "sqlite",
        # Images
        "photo.png": "image",
        "photo.jpg": "image",
        "photo.jpeg": "image",
        "anim.gif": "image",
        "modern.webp": "image",
        "old.bmp": "image",
        "scan.tiff": "image",
        "scan.tif": "image",
        "diagram.svg": "svg",
        # Audio
        "song.mp3": "audio",
        "clip.wav": "audio",
        "voice.m4a": "audio",
        "track.ogg": "audio",
        "lossless.flac": "audio",
        "stream.aac": "audio",
        # Video
        "movie.mp4": "video",
        "clip.mov": "video",
        "web.webm": "video",
        "old.avi": "video",
        "hd.mkv": "video",
        # Archives
        "files.zip": "zip",
        "backup.tar": "tar",
        "backup.gz": "tar",
        "backup.bz2": "tar",
        "archive.rar": "rar",
        # Email
        "mail.eml": "eml",
        "outlook.msg": "msg",
        # Web
        "page.html": "html",
        "page.htm": "html",
        # Code
        "script.py": "text",
        "app.js": "text",
        "app.ts": "text",
        "app.tsx": "text",
        "app.jsx": "text",
        "style.css": "text",
        "query.sql": "text",
        "run.sh": "text",
        "run.bash": "text",
        "run.ps1": "text",
        "run.bat": "text",
        "lib.rs": "text",
        "main.go": "text",
        "Main.java": "text",
        "lib.c": "text",
        "lib.cpp": "text",
        "lib.h": "text",
        "app.rb": "text",
        "index.php": "text",
        "App.swift": "text",
        "App.kt": "text",
        "analysis.r": "text",
        "game.lua": "text",
        "config.ini": "text",
        "config.cfg": "text",
        "secrets.env": "text",
        "debug.log": "text",
    }

    for filename, expected in cases.items():
        detected = detect_handler(filename)
        result(filename, detected == expected,
               f"expected={expected}, got={detected}" if detected != expected else "")

    # Test MIME type override
    detected = detect_handler("unknown_file", content_type="application/pdf")
    result("MIME override (application/pdf)", detected == "pdf")

    # Test unknown file
    detected = detect_handler("mystery.xyz")
    result("Unknown extension -> 'unknown'", detected == "unknown")


# ============================================
# PART 2: Process Files (end-to-end)
# ============================================

async def test_process_text():
    print("\n--- Text Files ---")
    data = b"Hello, CHAMP! This is a test file."
    r = await process_file(data, "test.txt")
    result("TXT", r.file_type == "text" and "Hello" in r.text, f"type={r.file_type}")

    py_data = b"def hello():\n    print('Hello CHAMP')\n"
    r = await process_file(py_data, "test.py")
    result("PY (code)", r.file_type == "text" and "python" in r.text, f"lang={r.metadata.get('language')}")

    js_data = b"console.log('hello');\n"
    r = await process_file(js_data, "app.js")
    result("JS (code)", r.file_type == "text" and "javascript" in r.text)

    md_data = b"# Title\n\nSome **bold** text.\n"
    r = await process_file(md_data, "readme.md")
    result("MD", r.file_type == "text" and "Title" in r.text)


async def test_process_csv():
    print("\n--- CSV/TSV ---")
    csv_data = b"name,age,city\nAlice,30,NYC\nBob,25,LA\n"
    r = await process_file(csv_data, "data.csv")
    result("CSV", r.file_type == "csv" and "Alice" in r.text and r.metadata.get("rows") == 3,
           f"rows={r.metadata.get('rows')}")

    tsv_data = b"name\tage\nAlice\t30\n"
    r = await process_file(tsv_data, "data.tsv")
    result("TSV", r.file_type == "csv" and "Alice" in r.text)


async def test_process_json():
    print("\n--- JSON ---")
    data = json.dumps({"name": "CHAMP", "version": 3, "features": ["voice", "memory"]}).encode()
    r = await process_file(data, "config.json")
    result("JSON", r.file_type == "json" and "CHAMP" in r.text)


async def test_process_xml():
    print("\n--- XML ---")
    data = b'<?xml version="1.0"?><config><name>CHAMP</name><version>3</version></config>'
    r = await process_file(data, "config.xml")
    result("XML", r.file_type == "xml" and "CHAMP" in r.text)


async def test_process_yaml():
    print("\n--- YAML ---")
    data = b"name: CHAMP\nversion: 3\nfeatures:\n  - voice\n  - memory\n"
    r = await process_file(data, "config.yaml")
    result("YAML", r.file_type == "yaml" and "CHAMP" in r.text)


async def test_process_toml():
    print("\n--- TOML ---")
    data = b'[project]\nname = "CHAMP"\nversion = 3\n'
    r = await process_file(data, "config.toml")
    result("TOML", r.file_type == "toml" and "CHAMP" in r.text)


async def test_process_html():
    print("\n--- HTML ---")
    data = b"<html><head><title>CHAMP</title></head><body><h1>Hello</h1><p>World</p></body></html>"
    r = await process_file(data, "page.html")
    result("HTML", r.file_type == "html" and "Hello" in r.text and r.metadata.get("title") == "CHAMP",
           f"title={r.metadata.get('title')}")


async def test_process_svg():
    print("\n--- SVG ---")
    data = b'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100"><circle cx="50" cy="50" r="40"/></svg>'
    r = await process_file(data, "diagram.svg")
    result("SVG (as XML text)", r.file_type == "svg" and "circle" in r.text)


async def test_process_image():
    print("\n--- Images ---")
    from PIL import Image

    # Create a small test PNG
    img = Image.new("RGB", (10, 10), color="red")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    png_data = buf.getvalue()

    r = await process_file(png_data, "test.png")
    result("PNG", r.file_type == "image" and r.image_b64 is not None and r.metadata.get("width") == 10,
           f"{r.metadata.get('width')}x{r.metadata.get('height')}")

    # JPEG
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    jpg_data = buf.getvalue()
    r = await process_file(jpg_data, "test.jpg")
    result("JPG", r.file_type == "image" and r.image_b64 is not None)

    # GIF
    buf = io.BytesIO()
    img.save(buf, format="GIF")
    gif_data = buf.getvalue()
    r = await process_file(gif_data, "test.gif")
    result("GIF", r.file_type == "image" and r.image_b64 is not None)

    # BMP
    buf = io.BytesIO()
    img.save(buf, format="BMP")
    bmp_data = buf.getvalue()
    r = await process_file(bmp_data, "test.bmp")
    result("BMP", r.file_type == "image" and r.image_b64 is not None)

    # WEBP
    buf = io.BytesIO()
    img.save(buf, format="WEBP")
    webp_data = buf.getvalue()
    r = await process_file(webp_data, "test.webp")
    result("WEBP", r.file_type == "image" and r.image_b64 is not None)

    # TIFF
    buf = io.BytesIO()
    img.save(buf, format="TIFF")
    tiff_data = buf.getvalue()
    r = await process_file(tiff_data, "test.tiff")
    result("TIFF", r.file_type == "image" and r.image_b64 is not None)


async def test_process_zip():
    print("\n--- Archives ---")
    # Create in-memory ZIP
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("hello.txt", "Hello CHAMP")
        zf.writestr("data.csv", "a,b\n1,2\n")
    zip_data = buf.getvalue()

    r = await process_file(zip_data, "test.zip")
    result("ZIP", r.file_type == "zip" and r.metadata.get("file_count") == 2 and "hello.txt" in r.text,
           f"files={r.metadata.get('file_count')}")


async def test_process_tar():
    import tarfile as tf
    buf = io.BytesIO()
    with tf.open(fileobj=buf, mode="w") as tar:
        info = tf.TarInfo(name="hello.txt")
        content = b"Hello CHAMP"
        info.size = len(content)
        tar.addfile(info, io.BytesIO(content))
    tar_data = buf.getvalue()

    r = await process_file(tar_data, "test.tar")
    result("TAR", r.file_type == "tar" and "hello.txt" in r.text)


async def test_process_eml():
    print("\n--- Email ---")
    eml_data = (
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Subject: Test Email for CHAMP\r\n"
        b"Date: Fri, 14 Mar 2026 12:00:00 +0000\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Hello Bob, this is a test email.\r\n"
    )
    r = await process_file(eml_data, "test.eml")
    result("EML", r.file_type == "eml" and "alice" in r.text and "Test Email" in r.text,
           f"subject={r.metadata.get('subject')}")


async def test_process_rtf():
    print("\n--- Rich Documents ---")
    rtf_data = br'{\rtf1\ansi Hello from CHAMP RTF}'
    r = await process_file(rtf_data, "test.rtf")
    result("RTF", r.file_type == "rtf" and "Hello" in r.text)


async def test_process_docx():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from CHAMP DOCX test")
    doc.add_paragraph("Second paragraph here")
    buf = io.BytesIO()
    doc.save(buf)
    docx_data = buf.getvalue()

    r = await process_file(docx_data, "test.docx")
    result("DOCX", r.file_type == "docx" and "Hello" in r.text and r.metadata.get("paragraphs", 0) >= 2,
           f"paragraphs={r.metadata.get('paragraphs')}")


async def test_process_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "CHAMP V3"
    slide.placeholders[1].text = "File Processor Test"
    buf = io.BytesIO()
    prs.save(buf)
    pptx_data = buf.getvalue()

    r = await process_file(pptx_data, "test.pptx")
    result("PPTX", r.file_type == "pptx" and "CHAMP" in r.text,
           f"slides={r.metadata.get('slides')}")


async def test_process_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Name", "Score"])
    ws.append(["CHAMP", 100])
    ws.append(["Test", 99])
    buf = io.BytesIO()
    wb.save(buf)
    xlsx_data = buf.getvalue()

    r = await process_file(xlsx_data, "test.xlsx")
    result("XLSX", r.file_type == "xlsx" and "CHAMP" in r.text and "100" in r.text,
           f"sheets={r.metadata.get('sheets')}")


async def test_process_pdf():
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello from CHAMP PDF test")
    pdf_data = doc.tobytes()
    doc.close()

    r = await process_file(pdf_data, "test.pdf")
    result("PDF", r.file_type == "pdf" and "Hello" in r.text,
           f"pages={r.metadata.get('pages')}")


async def test_process_audio():
    print("\n--- Audio (requires ffmpeg) ---")
    try:
        from pydub import AudioSegment
        from pydub.generators import Sine

        # Generate 1-second sine wave
        tone = Sine(440).to_audio_segment(duration=1000)

        # WAV
        buf = io.BytesIO()
        tone.export(buf, format="wav")
        wav_data = buf.getvalue()
        r = await process_file(wav_data, "test.wav")
        result("WAV", r.file_type == "audio" and r.metadata.get("duration_seconds", 0) > 0,
               f"duration={r.metadata.get('duration_seconds', 0):.1f}s")

        # MP3
        buf = io.BytesIO()
        tone.export(buf, format="mp3")
        mp3_data = buf.getvalue()
        r = await process_file(mp3_data, "test.mp3")
        result("MP3", r.file_type == "audio" and r.metadata.get("duration_seconds", 0) > 0,
               f"duration={r.metadata.get('duration_seconds', 0):.1f}s")

        # OGG
        buf = io.BytesIO()
        tone.export(buf, format="ogg")
        ogg_data = buf.getvalue()
        r = await process_file(ogg_data, "test.ogg")
        result("OGG", r.file_type == "audio" and r.metadata.get("duration_seconds", 0) > 0,
               f"duration={r.metadata.get('duration_seconds', 0):.1f}s")

        # FLAC
        buf = io.BytesIO()
        tone.export(buf, format="flac")
        flac_data = buf.getvalue()
        r = await process_file(flac_data, "test.flac")
        result("FLAC", r.file_type == "audio" and r.metadata.get("duration_seconds", 0) > 0,
               f"duration={r.metadata.get('duration_seconds', 0):.1f}s")

    except Exception as e:
        skip("Audio formats", f"ffmpeg not available or pydub error: {e}")


async def test_process_video():
    print("\n--- Video (requires ffmpeg) ---")
    try:
        import numpy as np

        # Create a tiny MP4 using moviepy
        try:
            from moviepy import VideoClip
        except ImportError:
            from moviepy.editor import VideoClip

        def make_frame(t):
            # 10x10 red frame
            return np.full((10, 10, 3), [255, 0, 0], dtype=np.uint8)

        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as f:
            tmp_path = f.name

        clip = VideoClip(make_frame, duration=1)
        clip.write_videofile(tmp_path, fps=10, codec="libx264",
                           audio=False, logger=None,
                           preset="ultrafast", threads=1)
        clip.close()

        with open(tmp_path, "rb") as f:
            mp4_data = f.read()
        os.unlink(tmp_path)

        r = await process_file(mp4_data, "test.mp4")
        result("MP4", r.file_type == "video" and r.metadata.get("duration_seconds", 0) > 0,
               f"duration={r.metadata.get('duration_seconds', 0):.1f}s, "
               f"res={r.metadata.get('width')}x{r.metadata.get('height')}")

    except Exception as e:
        skip("Video formats", f"ffmpeg/moviepy error: {e}")


async def test_process_sqlite():
    print("\n--- Database ---")
    import sqlite3

    with tempfile.NamedTemporaryFile(suffix=".db", delete=False) as f:
        tmp_path = f.name

    conn = sqlite3.connect(tmp_path)
    conn.execute("CREATE TABLE users (id INTEGER, name TEXT)")
    conn.execute("INSERT INTO users VALUES (1, 'CHAMP')")
    conn.execute("INSERT INTO users VALUES (2, 'Libby')")
    conn.commit()
    conn.close()

    with open(tmp_path, "rb") as f:
        db_data = f.read()
    os.unlink(tmp_path)

    r = await process_file(db_data, "test.db")
    result("SQLite", r.file_type == "sqlite" and "CHAMP" in r.text and "users" in str(r.metadata.get("tables")),
           f"tables={r.metadata.get('tables')}")


async def test_process_parquet():
    print("\n--- Parquet ---")
    try:
        import pyarrow as pa
        import pyarrow.parquet as pq

        table = pa.table({"name": ["CHAMP", "Test"], "score": [100, 99]})
        buf = io.BytesIO()
        pq.write_table(table, buf)
        parquet_data = buf.getvalue()

        r = await process_file(parquet_data, "test.parquet")
        result("Parquet", r.file_type == "parquet" and "CHAMP" in r.text,
               f"cols={r.metadata.get('columns')}, rows={r.metadata.get('total_rows')}")
    except Exception as e:
        skip("Parquet", str(e))


async def test_process_unknown():
    print("\n--- Edge Cases ---")
    # Unknown binary (use bytes that fail UTF-8 decode to trigger binary path)
    r = await process_file(b"\x80\x81\x82\xff\xfe\x00\x01", "mystery.xyz")
    result("Unknown binary -> graceful fallback", r.file_type in ("binary", "text"),
           f"type={r.file_type}")

    # Unknown text
    r = await process_file(b"some random text content", "mystery.xyz2")
    # Actually .xyz2 should be unknown, let's test with something truly unknown
    r = await process_file(b"hello world", "file.zzz")
    result("Unknown text -> treated as text", r.file_type == "text" and "hello" in r.text)

    # Empty file
    r = await process_file(b"", "empty.txt")
    result("Empty file", r.file_type == "text")

    # to_dict works
    r = await process_file(b"test", "test.txt")
    d = r.to_dict()
    result("to_dict()", "text" in d and "file_type" in d and "filename" in d)


# ============================================
# MAIN
# ============================================

async def main():
    start = time.time()
    print("=" * 60)
    print("  CHAMP V3 — Gate Test: File Processor")
    print("  Testing all 30+ format handlers")
    print("=" * 60)

    # Part 1: Detection
    test_handler_detection()

    # Part 2: Processing
    print("\n=== PART 2: File Processing (end-to-end) ===")

    await test_process_text()
    await test_process_csv()
    await test_process_json()
    await test_process_xml()
    await test_process_yaml()
    await test_process_toml()
    await test_process_html()
    await test_process_svg()
    await test_process_image()
    await test_process_zip()
    await test_process_tar()
    await test_process_eml()
    await test_process_rtf()
    await test_process_docx()
    await test_process_pptx()
    await test_process_xlsx()
    await test_process_pdf()
    await test_process_sqlite()
    await test_process_parquet()
    await test_process_audio()
    await test_process_video()
    await test_process_unknown()

    elapsed = time.time() - start
    print("\n" + "=" * 60)
    print(f"  RESULTS: {PASS} passed / {FAIL} failed / {SKIP} skipped")
    print(f"  Time: {elapsed:.1f}s")
    print("=" * 60)

    if FAIL > 0:
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())